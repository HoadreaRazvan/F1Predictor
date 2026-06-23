import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = r"C:\Users\hoadr\Desktop\Licenta\f1-predictor"
PLOTS = os.path.join(ROOT, "outputs", "plots")
OUT = os.path.join(ROOT, "Prezentare_Licenta_Hoadrea.pptx")

WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK   = RGBColor(0x16, 0x16, 0x1D)
TXT   = RGBColor(0x33, 0x33, 0x3B)
INK2  = RGBColor(0x5E, 0x5E, 0x66)
MUTED = RGBColor(0x9A, 0x9A, 0xA0)
HAIR  = RGBColor(0xE6, 0xE6, 0xE2)
SURF  = RGBColor(0xF5, 0xF5, 0xF2)
SURF2 = RGBColor(0xEF, 0xEF, 0xEB)
RED   = RGBColor(0xE1, 0x06, 0x00)
REDTINT = RGBColor(0xFB, 0xEA, 0xEA)
GREEN = RGBColor(0x12, 0x80, 0x5A)
AMBER = RGBColor(0xC8, 0x87, 0x1A)
GRAYBAR = RGBColor(0xCB, 0xCB, 0xD0)

DISP  = "Segoe UI Semibold"
LIGHT = "Segoe UI Light"
SEMI  = "Segoe UI Semibold"
BODY  = "Segoe UI"
MONO  = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5
LM, RM = 0.85, 12.48
CW = RM - LM

def slide_new():
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = WHITE
    r.line.fill.background(); r.shadow.inherit = False
    return s

def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, radius=None):
    typ = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    sh = s.shapes.add_shape(typ, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        try: sh.adjustments[0] = radius
        except Exception: pass
    if fill is None: sh.fill.background()
    else: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None: sh.line.fill.background()
    else: sh.line.color.rgb = line; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh

def hairline(s, x, y, w, color=HAIR, th=0.013):
    return rect(s, x, y, w, th, fill=color)

def vline(s, x, y, h, color=HAIR, th=0.013):
    return rect(s, x, y, th, h, fill=color)

def tb(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    box = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    f = box.text_frame; f.word_wrap = True; f.vertical_anchor = anchor
    f.margin_left = 0; f.margin_right = 0; f.margin_top = 0; f.margin_bottom = 0
    return f

class Text:
    def __init__(self, tf): self.tf = tf; self.first = True
    def add(self, segs, align=PP_ALIGN.LEFT, sb=0, sa=6, ls=1.0):
        p = self.tf.paragraphs[0] if self.first else self.tf.add_paragraph()
        self.first = False
        p.alignment = align
        if sb: p.space_before = Pt(sb)
        p.space_after = Pt(sa); p.line_spacing = ls
        if isinstance(segs, str): segs = [(segs,)]
        for seg in segs:
            t = seg[0]
            size = seg[1] if len(seg) > 1 else 15
            color = seg[2] if len(seg) > 2 else TXT
            bold = seg[3] if len(seg) > 3 else False
            italic = seg[4] if len(seg) > 4 else False
            name = seg[5] if len(seg) > 5 else BODY
            r = p.add_run(); r.text = t
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.italic = italic; r.font.name = name
        return p

def header(s, kicker, title, idx, tsize=30):
    rect(s, LM, 0.6, 0.42, 0.12, fill=RED)
    Text(tb(s, LM + 0.55, 0.55, 9.0, 0.3)).add([(kicker.upper(), 12, RED, False, False, SEMI)], sa=0)
    Text(tb(s, LM, 0.86, 10.6, 0.75)).add([(title, tsize, INK, False, False, DISP)], sa=0)
    Text(tb(s, RM - 1.6, 0.5, 1.6, 0.7)).add(
        [(f"{idx:02d}", 30, RGBColor(0xEC, 0xEC, 0xE8), False, False, LIGHT)], align=PP_ALIGN.RIGHT, sa=0)

def footer(s, n):
    Text(tb(s, LM, 6.98, 8.0, 0.3)).add(
        [("Predicția podiumului în Formula 1  ·  R.-B. Hoadrea", 9, MUTED, False, False, BODY)], sa=0)
    Text(tb(s, RM - 1.4, 6.98, 1.4, 0.3)).add(
        [(f"{n:02d} / 16", 9, MUTED, False, False, SEMI)], align=PP_ALIGN.RIGHT, sa=0)

def notes(s, text): s.notes_slide.notes_text_frame.text = text

def points(s, x, y, w, h, items, size=15, gap=14, ls=1.12, body=TXT):
    T = Text(tb(s, x, y, w, h))
    for it in items:
        if isinstance(it, str):
            T.add([("—  ", size, RED, True), (it, size, body)], sa=gap, ls=ls)
        else:
            T.add([("—  ", size, RED, True)] + it, sa=gap, ls=ls)
    return T

def center_in(sh, segs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    f = sh.text_frame; f.word_wrap = True; f.vertical_anchor = anchor
    f.margin_left = Inches(0.06); f.margin_right = Inches(0.06)
    f.margin_top = Inches(0.02); f.margin_bottom = Inches(0.02)
    T = Text(f)
    if segs and isinstance(segs[0], tuple): T.add(segs, align=align, sa=0)
    else:
        for para in segs: T.add(para, align=align, sa=0)
    return sh

def figure(s, img, x, y, w, h, caption=None, frame=True, pad=0.12):
    if frame: rect(s, x, y, w, h, fill=WHITE, line=HAIR, line_w=1)
    path = os.path.join(PLOTS, img)
    iw, ih = Image.open(path).size; ratio = iw / ih
    aw, ah = w - 2*pad, h - 2*pad
    if aw/ah > ratio: dh = ah; dw = dh*ratio
    else: dw = aw; dh = dw/ratio
    s.shapes.add_picture(path, Inches(x+(w-dw)/2), Inches(y+(h-dh)/2), Inches(dw), Inches(dh))
    if caption:
        Text(tb(s, x, y+h+0.06, w, 0.3)).add(
            [(caption, 10.5, MUTED, False, True, BODY)], align=PP_ALIGN.CENTER, sa=0)

def etable(s, x, y, headers, rows, col_w, aligns, row_h=0.5, hsize=11, bsize=13):
    total = sum(col_w)
    def ppal(a): return PP_ALIGN.LEFT if a == 'l' else (PP_ALIGN.RIGHT if a == 'r' else PP_ALIGN.CENTER)
    cx = x
    for j, h in enumerate(headers):
        Text(tb(s, cx, y, col_w[j], 0.32)).add([(h.upper(), hsize, RED, False, False, SEMI)], align=ppal(aligns[j]), sa=0)
        cx += col_w[j]
    cy = y + 0.38
    hairline(s, x, cy - 0.04, total, INK if False else HAIR)
    for row in rows:
        cx = x
        for j, cell in enumerate(row):
            if isinstance(cell, tuple): txt, color, bold = cell
            else: txt, color, bold = cell, (INK if j == 0 else TXT), (j == 0)
            fnt = BODY if aligns[j] == 'l' else SEMI
            Text(tb(s, cx, cy + 0.05, col_w[j], row_h - 0.05, anchor=MSO_ANCHOR.MIDDLE)).add(
                [(str(txt), bsize, color, bold, False, fnt)], align=ppal(aligns[j]), sa=0)
            cx += col_w[j]
        cy += row_h
        hairline(s, x, cy - 0.04, total, HAIR)
    return cy

def stat(s, x, y, w, num, lab, nsize=38, ncolor=INK):
    Text(tb(s, x, y, w, 0.66)).add([(num, nsize, ncolor, False, False, LIGHT)], sa=0)
    Text(tb(s, x, y + 0.66, w, 0.4)).add([(lab, 12, MUTED, False, False, BODY)], ls=1.0, sa=0)

s = slide_new()
Text(tb(s, LM, 1.05, 9.0, 0.3)).add([("PROIECT DE DIPLOMĂ  —  SIBIU 2026", 13, RED, False, False, SEMI)], sa=0)
rect(s, LM, 1.45, 1.25, 0.05, fill=RED)
Text(tb(s, LM, 1.75, 8.5, 2.4)).add(
    [("Comparare algoritmi de predicție cu învățare automată", 41, INK, False, False, LIGHT)], ls=1.04, sa=0)
Text(tb(s, LM, 4.05, 8.0, 1.0)).add(
    [("Predicția podiumului în Formula 1 cu trei algoritmi de învățare automată ", 17, INK2),
     ("implementați de la zero", 17, INK, False, False, SEMI)], ls=1.15, sa=0)
hairline(s, LM, 5.25, 7.7)
ab = Text(tb(s, LM, 5.5, 8.2, 1.6))
ab.add([("Absolvent   ", 12, MUTED), ("Răzvan-Bucur Hoadrea", 14, INK, False, False, SEMI),
        ("  ·  Specializarea Calculatoare", 13, INK2)], sa=8)
ab.add([("Îndrumător   ", 12, MUTED), ("Șef lucr. Alexandru Dorobanțiu", 14, INK, False, False, SEMI)], sa=8)
ab.add([("Universitatea „Lucian Blaga” din Sibiu  ·  Facultatea de Inginerie", 12, INK2)], sa=0)

base = 5.0
for px, ph, pc, rk in [(10.05, 1.3, GRAYBAR, "2"), (10.93, 2.05, RED, "1"), (11.81, 0.95, RGBColor(0xE2,0xE2,0xE5), "3")]:
    rect(s, px, base - ph, 0.66, ph, fill=pc, radius=0.08)
    Text(tb(s, px, base - ph - 0.36, 0.66, 0.3)).add(
        [(rk, 15, INK, False, False, SEMI)], align=PP_ALIGN.CENTER, sa=0)
hairline(s, 9.85, base, 2.75, color=INK, th=0.02)
Text(tb(s, 9.85, base + 0.12, 2.75, 0.3)).add([("P1   ·   P2   ·   P3", 11, MUTED, False, False, SEMI)], align=PP_ALIGN.CENTER, sa=0)
notes(s, "Bună ziua. Numele meu este Răzvan-Bucur Hoadrea și vă prezint proiectul de diplomă "
         "„Comparare algoritmi de predicție cu învățare automată”, îndrumat de domnul Șef lucrări "
         "Alexandru Dorobanțiu. Lucrarea prezice podiumul curselor de Formula 1 folosind trei "
         "algoritmi de învățare automată pe care i-am implementat complet de la zero. În următoarele "
         "minute voi parcurge problema, soluția și, mai ales, studiul comparativ dintre cei trei algoritmi.")

s = slide_new()
header(s, "Introducere", "Tema și contextul", 2)
points(s, LM, 1.95, 7.2, 4.6, [
    "Predicția rezultatelor sportive este un teren clasic pentru învățarea automată: date structurate, abundente, actualizate periodic.",
    "Formula 1 — cadru ideal: zeci de curse pe sezon, ~20 de piloți pe cursă, reguli stabile pe mai mulți ani.",
    [("Problema lucrării: ", 15, INK, True), ("prezicerea podiumului (P1, P2, P3) al unei curse.", 15, TXT)],
    [("Dublă predicție: ", 15, INK, True), ("podium per cursă + clasament de campionat estimat.", 15, TXT)],
], size=15, gap=16)
vline(s, 9.05, 2.0, 4.0)
sx = 9.45
facts = [("~20", "piloți pe cursă"), ("3", "locuri de podium"), ("≈15%", "exemple pozitive (podium)")]
yy = 2.05
for big, lab in facts:
    Text(tb(s, sx, yy, 3.0, 0.6)).add([(big, 33, INK, False, False, LIGHT)], sa=0)
    Text(tb(s, sx, yy + 0.62, 3.0, 0.4)).add([(lab, 12.5, INK2)], sa=0)
    yy += 1.32
footer(s, 2)
notes(s, "Predicția sportivă e un domeniu clasic de ML. Formula 1 e deosebit de potrivită: fiecare "
         "sezon are zeci de curse, fiecare cursă ~20 de piloți, iar regulamentul e stabil. Spre "
         "deosebire de sporturile de echipă, performanța individuală se poate urmări numeric. "
         "Problema pe care o rezolv este prezicerea podiumului — primii trei piloți — al unei curse, "
         "iar din probabilități construiesc și un clasament de campionat. Rețineți cifra cheie: doar "
         "circa 15% dintre exemple sunt pozitive, un dezechilibru care va reveni constant în discuție.")

s = slide_new()
header(s, "De ce „de la zero”", "Motivație și obiective", 3)
Text(tb(s, LM, 1.95, 5.4, 0.3)).add([("MOTIVAȚIE", 12, MUTED, False, False, SEMI)], sa=0)
points(s, LM, 2.35, 5.5, 4.2, [
    [("Implementarea ", 14, TXT), ("de la zero", 14, INK, True), (" expune mecanismele ascunse de scikit-learn.", 14, TXT)],
    [("Scurgerea de informație", 14, INK, True), (" (data leakage) — o provocare metodologică reală, cauzalitate strictă.", 14, TXT)],
    [("Date dezechilibrate", 14, INK, True), (" (~15% podium) → acuratețea globală e înșelătoare.", 14, TXT)],
], size=14, gap=14)
vline(s, 6.75, 2.0, 4.55)
Text(tb(s, 7.1, 1.95, 5.3, 0.3)).add([("OBIECTIVE", 12, MUTED, False, False, SEMI)], sa=0)
objs = [
    "Implementarea de la zero a 3 modele (logreg, arbore, pădure), doar cu NumPy.",
    "Aplicare pe date reale F1 (2018–2025), cu feature engineering strict cauzal.",
    "Studiu comparativ pe 4 metrici de calitate a predicției.",
    "Evaluare cu separare temporală train / validare / test pe ani.",
    "Analiza efectului datelor in-sezon (antrenare incrementală).",
]
oy = 2.4
for i, o in enumerate(objs, 1):
    Text(tb(s, 7.1, oy - 0.04, 0.55, 0.5)).add([(str(i), 22, RED, False, False, LIGHT)], sa=0)
    Text(tb(s, 7.7, oy + 0.02, 4.7, 0.6, anchor=MSO_ANCHOR.MIDDLE)).add([(o, 13.5, TXT)], ls=1.04, sa=0)
    oy += 0.83
footer(s, 3)
notes(s, "Tema e motivată de trei lucruri. Întâi, scriind algoritmii de la zero înțeleg în "
         "profunzime ce ascund bibliotecile gata făcute. Al doilea, datele de F1 ridică o problemă "
         "metodologică reală — scurgerea de informație: orice caracteristică trebuie calculată doar "
         "din trecut. Al treilea, clasele sunt dezechilibrate, deci acuratețea simplă e înșelătoare. "
         "Pornind de aici am formulat cinci obiective concrete, pe care le veți regăsi îndeplinite "
         "la finalul prezentării.")

s = slide_new()
header(s, "Clasificare → ordonare", "Formularea problemei", 4)
points(s, LM, 1.95, 11.5, 2.0, [
    [("Pentru fiecare pereche ", 15, TXT), ("(cursă, pilot)", 15, INK, True, False, MONO),
     (" definim eticheta ", 15, TXT), ("is_podium ∈ {0, 1}", 15, INK, True, False, MONO), ("  (1 dacă locul ≤ 3).", 15, TXT)],
    [("Un clasificator binar estimează probabilitatea ", 15, TXT), ("P(podium) ∈ [0, 1]", 15, INK, True, False, MONO), (".", 15, TXT)],
    [("Pentru fiecare cursă, sortăm piloții ", 15, TXT), ("descrescător după P(podium)", 15, INK, True),
     (" → primii 3 = podiumul prezis.", 15, TXT)],
], size=15, gap=12)

fy = 4.25; chips = ["(cursă, pilot)", "P(podium)", "sortare desc.", "Top-3 = podium"]
chw = 2.55; aw = 0.55; fx = LM
for i, c in enumerate(chips):
    fill = REDTINT if i == len(chips)-1 else SURF
    txtc = RED if i == len(chips)-1 else INK
    chip = rect(s, fx, fy, chw, 0.92, fill=fill, radius=0.06)
    center_in(chip, [[(c, 14.5, txtc, True, False, (MONO if i < 2 else SEMI))]])
    if i < len(chips)-1:
        Text(tb(s, fx+chw, fy, aw, 0.92, anchor=MSO_ANCHOR.MIDDLE)).add(
            [("→", 22, RED, False, False, BODY)], align=PP_ALIGN.CENTER, sa=0)
    fx += chw + aw
Text(tb(s, LM, 5.7, 11.4, 0.9)).add(
    [("De ce ordonare (ranking)?  ", 14, RED, True),
     ("un clasificator binar aplicat independent pe fiecare pilot nu garantează exact 3 predicții "
      "pozitive pe cursă; ordonarea produce întotdeauna un podium bine definit.", 14, TXT)], ls=1.1, sa=0)
footer(s, 4)
notes(s, "Iată cum am formulat problema. Pentru fiecare pereche cursă-pilot pun eticheta is_podium, "
         "1 dacă pilotul termină în primii trei. Modelul nu dă direct eticheta, ci o probabilitate de "
         "podium. Cheia e pasul următor: pentru o cursă sortez toți piloții după această probabilitate "
         "și iau primii trei ca podium prezis, cel mai probabil fiind P1. Practic transform clasificarea "
         "într-o problemă de ordonare. E important pentru că un clasificator binar independent nu îmi "
         "garantează fix trei piloți pe podium — ordonarea, da.")

s = slide_new()
header(s, "Soluția", "Cei trei algoritmi — implementați de la zero peste NumPy", 5, tsize=26)
cards = [
    ("01", "Regresie logistică", ["Model liniar + sigmoid", "Gradient descent pe cross-entropy",
        "Regularizare L2, ponderare pe clase"], "bias mare · varianță mică", "σ(z) = 1 / (1 + e⁻ᶻ)"),
    ("02", "Arbore de decizie (CART)", ["Partiționare recursivă a spațiului", "Criteriu: entropie + information gain",
        "Neliniar, interpretabil"], "bias mic · varianță MARE", "IG = H(p) − Σ (nₖ/n)·H(k)"),
    ("03", "Random Forest", ["Ansamblu de arbori de decizie", "Bagging + subspații aleatoare de feature-uri",
        "Mediere → decorelare"], "bias mic · varianță mică", "p̂ = medie( arbori )"),
]
cx = LM; cw = 3.78; gap = 0.14; cy = 2.0; ch = 3.95
for idx, name, items, tag, formula in cards:
    rect(s, cx, cy, cw, ch, fill=SURF, radius=0.035)
    Text(tb(s, cx + 0.25, cy + 0.22, 1.2, 0.4)).add([(idx, 15, RED, False, False, SEMI)], sa=0)
    Text(tb(s, cx + 0.25, cy + 0.55, cw - 0.5, 0.7)).add([(name, 15.5, INK, False, False, DISP)], ls=1.0, sa=0)
    T = Text(tb(s, cx + 0.25, cy + 1.35, cw - 0.5, 1.6))
    for it in items:
        T.add([("—  ", 12.5, RED, True), (it, 12.5, TXT)], sa=9, ls=1.04)
    Text(tb(s, cx + 0.25, cy + ch - 1.0, cw - 0.5, 0.3)).add([(tag, 12, INK, True, True, BODY)], sa=0)
    chip = rect(s, cx + 0.25, cy + ch - 0.6, cw - 0.5, 0.4, fill=WHITE, radius=0.1)
    center_in(chip, [[(formula, 12.5, INK, False, False, MONO)]])
    cx += cw + gap
Text(tb(s, LM, 6.15, 11.5, 0.6)).add(
    [("Spectrul bias–varianță:  ", 13.5, RED, True),
     ("pădurea reduce varianța arborelui prin mediere, păstrând bias-ul scăzut.", 13.5, TXT)], sa=0)
footer(s, 5)
notes(s, "Am implementat trei algoritmi, fără nicio bibliotecă de ML — doar NumPy pentru algebră. "
         "Regresia logistică: model liniar trecut prin sigmoid, antrenat cu gradient descent pe "
         "cross-entropy, cu regularizare L2 — bias mare, varianță mică. Arborele de decizie CART "
         "partiționează recursiv spațiul după information gain bazat pe entropie; e flexibil, dar are "
         "varianță mare. Random Forest combină mulți arbori prin bagging și subspații aleatoare de "
         "caracteristici, iar medierea reduce varianța. Cele trei acoperă intenționat tot spectrul "
         "bias–varianță, ceea ce face comparația interesantă.")

s = slide_new()
header(s, "Sursa de date", "Datele — FastF1 (2018–2025)", 6)
stats = [("8", "sezoane"), ("173", "curse"), ("3458", "rezultate"), ("43", "piloți"), ("19", "echipe")]
sx = LM; colw = CW / 5
for i, (big, lab) in enumerate(stats):
    Text(tb(s, sx + 0.05, 1.95, colw - 0.2, 0.7)).add([(big, 37, INK, False, False, LIGHT)], sa=0)
    Text(tb(s, sx + 0.07, 2.66, colw - 0.2, 0.4)).add([(lab, 12.5, INK2)], sa=0)
    if i < 4: vline(s, sx + colw - 0.1, 2.0, 0.95)
    sx += colw
hairline(s, LM, 3.35, CW)
points(s, LM, 3.6, 11.5, 3.0, [
    [("Sursa: ", 15, INK, True), ("FastF1", 15, TXT), (" (rezultate oficiale, calificări, vreme).  ", 15, TXT),
     ("Acces strict offline", 15, INK, True), (" → reproductibil, fără internet.", 15, TXT)],
    [("Per cursă: ", 15, INK, True), ("sesiunea R (grilă, poziție, status, puncte, echipă) + sesiunea Q (calificări) + vreme.", 15, TXT)],
    [("Îmbogățiri: ", 15, INK, True), ("timp la boxe, ritm de cursă, apariția Safety Car, degradarea pneurilor.", 15, TXT)],
    [("Dezechilibru: ", 15, INK, True), ("rata de podium ≈ 0,15 — un model „niciun podium” ar avea deja ~85% acuratețe.", 15, TXT)],
], size=15, gap=13)
footer(s, 6)
notes(s, "Datele vin din biblioteca FastF1: opt sezoane consecutive, 2018–2025, însemnând 173 de "
         "curse și 3458 de rezultate cursă-pilot, de la 43 de piloți și 19 echipe. Un detaliu de "
         "proiectare important: totul rulează strict offline, dintr-un cache local, deci rezultatele "
         "sunt complet reproductibile. Pentru fiecare cursă încarc sesiunea de cursă și cea de "
         "calificări, plus vremea, și adaug semnale îmbogățite precum ritmul de cursă sau degradarea "
         "pneurilor. Subliniez din nou dezechilibrul: rata de podium e 0,15.")

s = slide_new()
header(s, "Etapa centrală a soluției", "Feature engineering fără data leakage", 7, tsize=28)
Text(tb(s, LM, 1.95, 6.7, 1.0)).add(
    [("44 de caracteristici numerice, toate ", 19, INK, False, False, LIGHT),
     ("strict cauzale", 19, RED, False, False, DISP), (".", 19, INK, False, False, LIGHT)], ls=1.08, sa=0)
points(s, LM, 3.0, 6.8, 3.5, [
    [("Principiul cauzal: ", 14.5, INK, True), ("feature-urile pentru cursa R folosesc DOAR rezultate din curse strict anterioare.", 14.5, TXT)],
    [("Procesare ", 14.5, TXT), ("cronologică", 14.5, INK, True), ("; starea istorică se actualizează ", 14.5, TXT),
     ("după", 14.5, INK, True), (" calcularea feature-urilor.", 14.5, TXT)],
    [("Standardizarea (μ, σ) ", 14.5, TXT), ("estimată doar pe train", 14.5, INK, True), (", reaplicată identic la validare/test.", 14.5, TXT)],
    [("Garanție verificată ", 14.5, TXT), ("automat", 14.5, INK, True), (" printr-un test anti-leakage.", 14.5, TXT)],
], size=14.5, gap=13)
vline(s, 8.15, 2.0, 4.5)
cat = [["Grilă / calificări", "9"], ["Formă rulantă pilot", "8"], ["Sezon-to-date pilot", "5"],
       ["Forța echipei", "8"], ["Istoric circuit", "8"], ["Context (experiență, vreme)", "6"],
       [("Total", INK, True), ("44", GREEN, True)]]
etable(s, 8.5, 2.05, ["Categorie", "Nr."], cat, [3.2, 0.75], ['l', 'r'], row_h=0.52)
footer(s, 7)
notes(s, "Aceasta e, în opinia mea, cea mai importantă parte metodologică. Am construit 44 de "
         "caracteristici, toate strict cauzale. Principiul: orice caracteristică pentru cursa R se "
         "calculează exclusiv din curse strict anterioare. Procesez cursele în ordine cronologică și "
         "actualizez starea istorică abia DUPĂ ce am calculat caracteristicile. Chiar și standardizarea "
         "folosește media și deviația estimate doar pe train. Iar absența scurgerii e verificată printr-un "
         "test automat. Caracteristicile se grupează în șase categorii — poziția de start, forma "
         "pilotului, forța echipei, istoricul pe circuit și contextul.")

s = slide_new()
header(s, "Pipeline modular", "Arhitectura sistemului", 8)
stages = ["Cache FastF1", "Tabel de\nrezultate", "Caracteristici\ncauzale (44)",
          "Modele\nlogreg·tree·forest", "Predicție\n(probabilități)", "Evaluare\n& grafice"]
n = len(stages); bw = 1.66; aw = 0.3; bh = 1.15; by = 2.35
total = n*bw + (n-1)*aw; bx = (SW - total)/2
xs = []
for i, lab in enumerate(stages):
    fill = REDTINT if i in (3, 4) else SURF
    box = rect(s, bx, by, bw, bh, fill=fill, radius=0.05)
    center_in(box, [[(line, 11.5, (RED if i in (3, 4) else INK), True, False, BODY)] for line in lab.split("\n")])
    xs.append(bx)
    if i < n-1:
        Text(tb(s, bx+bw, by, aw, bh, anchor=MSO_ANCHOR.MIDDLE)).add([("→", 17, RED, False, False, BODY)], align=PP_ALIGN.CENTER, sa=0)
    bx += bw + aw
oy = by + bh + 0.5
Text(tb(s, xs[3], oy - 0.3, 5.0, 0.3)).add([("ieșiri", 11, MUTED, False, True)], sa=0)
for j, lab in enumerate(["Podium per cursă (Top-3)", "Clasament de campionat"]):
    ob = rect(s, xs[3] + j*2.55, oy, 2.35, 0.6, fill=WHITE, line=HAIR, line_w=1, radius=0.06)
    center_in(ob, [[(lab, 11.5, INK, True, False, BODY)]])
points(s, LM, 5.65, 11.5, 1.2, [
    [("Modular și testabil pe etape", 14, INK, True), (" — corectitudinea cauzală se verifică izolat, fără a antrena vreun model.", 14, TXT)],
    [("Etapele costisitoare sunt memorate pe disc (Parquet)", 14, INK, True), (" și reutilizate; ușor de extins cu un nou model.", 14, TXT)],
], size=14, gap=9)
footer(s, 8)
notes(s, "Sistemul e un pipeline modular. Pornesc de la cache-ul FastF1, construiesc un tabel de "
         "rezultate, apoi îl îmbogățesc cu cele 44 de caracteristici cauzale. Pe acest set antrenez "
         "cele trei modele, ale căror probabilități alimentează atât predicția podiumului per cursă, "
         "cât și clasamentul de campionat, plus modulul de evaluare și grafice. Fiecare etapă e un "
         "subpachet separat, testabil independent, iar etapele scumpe sunt memorate pe disc.")

s = slide_new()
header(s, "Fără scurgere de informație", "Validare temporală pe ani", 9)
years = list(range(2018, 2026)); cwc = 1.28; cg = 0.13; total = len(years)*cwc + (len(years)-1)*cg
startx = (SW - total)/2
xx = startx
for y in years:
    box = rect(s, xx, 1.95, cwc, 0.56, fill=SURF, radius=0.07)
    center_in(box, [[(str(y), 14, INK2, False, False, SEMI)]])
    xx += cwc + cg
Text(tb(s, startx, 2.72, total, 0.3)).add([("Exemplu — sezon de test 2024:", 13, INK, True)], sa=0)
xx = startx; rowy = 3.05
fold = {2018: INK, 2019: INK, 2020: INK, 2021: INK, 2022: INK, 2023: AMBER, 2024: RED, 2025: SURF2}
for y in years:
    col = fold[y]; tc = MUTED if y == 2025 else WHITE
    box = rect(s, xx, rowy, cwc, 0.56, fill=col, radius=0.07)
    center_in(box, [[(str(y), 14, tc, True, False, SEMI)]])
    xx += cwc + cg
lx = startx; ly = rowy + 0.82
for col, lab in [(INK, "Antrenare  (≤ S−2)"), (AMBER, "Validare  (S−1)"), (RED, "Test  (S)"), (SURF2, "neutilizat")]:
    rect(s, lx, ly, 0.3, 0.3, fill=col, radius=0.16)
    Text(tb(s, lx + 0.4, ly + 0.02, 2.5, 0.3)).add([(lab, 12.5, INK2)], sa=0)
    lx += 2.9
points(s, LM, 5.35, 11.5, 1.5, [
    [("Pentru fiecare sezon S: ", 14, INK, True), ("antrenare pe ≤ S−2, validare pe S−1, test pe S", 14, TXT),
     ("  →  prezicem mereu „viitorul” cu modele antrenate doar pe „trecut”.", 14, TXT)],
    [("6 sezoane de test (2020–2025), 131 de curse.", 14, INK, True),
     ("  Exemplele de antrenament cresc de la 420 (test 2020) la 2500 (test 2025).", 14, TXT)],
], size=14, gap=9)
footer(s, 9)
notes(s, "Evaluarea e punctul forte metodologic. Folosesc o validare temporală: pentru a prezice "
         "sezonul S, antrenez modelul doar pe sezoanele până la S minus 2, țin sezonul S minus 1 ca "
         "validare, și raportez performanța pe sezonul S, complet nevăzut. În exemplul de pe ecran, "
         "pentru testul pe 2024 antrenez pe 2018–2022, validez pe 2023 și testez pe 2024. Astfel niciun "
         "model nu vede vreodată sezonul pe care e evaluat. Am șase sezoane de test, 2020–2025, "
         "însemnând 131 de curse.")

s = slide_new()
header(s, "Cum măsurăm calitatea", "Metrici de evaluare", 10)
metrics = [
    ("01", "Acuratețe", "Proporția clasificărilor corecte (podium / non-podium)."),
    ("02", "Rată de nimerire a podiumului", "Câți dintre cei 3 piloți preziși se află în podiumul real."),
    ("03", "Podium exact", "Fracția curselor cu toți cei 3 piloți nimeriți (ca mulțime)."),
    ("04", "Câștigător (Top-1)", "Fracția curselor în care P1 este prezis corect."),
]
pos = [(LM, 1.95), (6.93, 1.95), (LM, 3.35), (6.93, 3.35)]
for (mx, my), (idx, mt, md) in zip(pos, metrics):
    rect(s, mx, my, 5.5, 1.25, fill=SURF, radius=0.04)
    Text(tb(s, mx + 0.28, my + 0.2, 1.0, 0.4)).add([(idx, 14, RED, False, False, SEMI)], sa=0)
    Text(tb(s, mx + 0.95, my + 0.18, 4.4, 0.45)).add([(mt, 15, INK, False, False, DISP)], sa=0)
    Text(tb(s, mx + 0.95, my + 0.6, 4.35, 0.6)).add([(md, 12.5, INK2)], ls=1.03, sa=0)
hairline(s, LM, 4.95, CW)
Text(tb(s, LM, 5.15, 11.5, 1.4)).add(
    [("Acuratețea este înșelătoare sub dezechilibru: ", 14.5, RED, True),
     ("un model care prezice mereu „niciun podium” atinge deja ~85%. De aceea metricile la nivel de "
      "cursă sunt cele decisive.", 14.5, TXT)], sa=8, ls=1.08)
Text(tb(s, LM, 6.0, 11.5, 0.5)).add(
    [("Metrici probabilistice: ", 14, INK, True), ("log-loss · Brier · ROC-AUC · PR-AUC · ECE.", 14, TXT),
     ("   Linie de bază aleatoare = pragul de șansă.", 14, INK2)], ls=1.05, sa=0)
footer(s, 10)
notes(s, "Pentru evaluare folosesc patru metrici principale, calculate la nivel de cursă și apoi "
         "mediate. Acuratețea clasificării binare; rata de nimerire a podiumului — câți din cei trei "
         "piloți preziși sunt corecți; podiumul exact — fracția curselor în care nimeresc toți trei; "
         "și acuratețea câștigătorului, adică P1. Atenție însă: din cauza dezechilibrului, acuratețea "
         "e înșelătoare — reperul trivial e deja 85%. De aceea metricile la nivel de cursă contează. "
         "Adaug și metrici probabilistice și compar mereu cu o linie de bază aleatoare.")

s = slide_new()
header(s, "Studiul comparativ", "Rezultate — performanța comparativă", 11, tsize=28)
rows = [
    ["Regresie logistică", "0,838", "0,679", "0,246", ("0,541", GREEN, True)],
    ["Arbore de decizie", "0,885", "0,640", "0,197", "0,293"],
    ["Random Forest", ("0,902", GREEN, True), ("0,688", GREEN, True), ("0,267", GREEN, True), "0,437"],
    [("Linie de bază", MUTED, False), ("0,527", MUTED, False), ("0,263", MUTED, False), ("0,000", MUTED, False), ("0,014", MUTED, False)],
]
etable(s, LM, 2.0, ["Model", "Acurat.", "Rată pod.", "Pod. exact", "Câștig."],
       rows, [2.4, 1.0, 1.05, 1.05, 1.0], ['l', 'c', 'c', 'c', 'c'], row_h=0.5, bsize=13)
figure(s, "metric_summary.png", 7.55, 1.95, 4.9, 2.7, caption="Cele 4 metrici, medii pe sezoanele de test")
points(s, LM, 4.55, 11.6, 2.2, [
    [("Random Forest domină 3 din 4 metrici", 14, INK, True), (" — avantajul metodelor de ansamblu pe date tabelare.", 14, TXT)],
    [("Regresia logistică e cea mai bună la câștigător (Top-1)", 14, INK, True), (" și se antrenează de ~10× mai repede.", 14, TXT)],
    [("Arborele unic e constant cel mai slab pe podium", 14, INK, True), (" — varianța mare a unui arbore singur.", 14, TXT)],
    [("Toate modelele depășesc net linia de bază aleatoare", 14, INK, True), (" → învață semnal real.", 14, TXT)],
], size=14, gap=8)
footer(s, 11)
notes(s, "Iată rezultatul central. În tabel sunt cele trei modele pe cele patru metrici, plus linia "
         "de bază aleatoare. Random Forest e cel mai bun pe trei din patru metrici — acuratețe 0,90, "
         "rată de podium 0,69 și podium exact 0,27 — confirmând avantajul ansamblurilor. Interesant, "
         "regresia logistică, după optimizarea hiperparametrilor, ia cea mai bună acuratețe a "
         "câștigătorului, 0,54, și e mult mai ieftină. Arborele singur rămâne cel mai slab pe podium, "
         "ilustrând varianța lui mare. Toate trei depășesc clar șansa.")

s = slide_new()
header(s, "Dincolo de ordine", "Probabilități și calibrare", 12)
rf = [("0,936", "ROC-AUC"), ("0,230", "log-loss"), ("0,023", "ECE")]
sx = LM
for big, lab in rf:
    Text(tb(s, sx, 1.95, 2.3, 0.6)).add([(big, 30, GREEN, False, False, LIGHT)], sa=0)
    Text(tb(s, sx, 2.58, 2.3, 0.3)).add([(lab, 12, INK2)], sa=0)
    sx += 2.4
Text(tb(s, LM, 2.95, 7.0, 0.3)).add([("Random Forest — cel mai bun pe toate metricile probabilistice", 12, MUTED, False, True)], sa=0)
points(s, LM, 3.4, 7.0, 3.2, [
    [("Logreg și arborele ordonează bine, dar emit probabilități ", 14, TXT), ("supraîncrezătoare", 14, INK, True), (" (log-loss mare).", 14, TXT)],
    [("Calibrare (Platt / izotonică): ", 14, INK, True), ("log-loss logreg ", 14, TXT), ("1,913 → 0,259", 14, GREEN, True),
     (", ECE ", 14, TXT), ("0,072 → 0,025", 14, GREEN, True), (".", 14, TXT)],
    [("Calibrarea ", 14, TXT), ("nu schimbă podiumul prezis", 14, INK, True), (" (e monotonă) — doar calitatea probabilităților.", 14, TXT)],
    [("Pădurea, mediind 60 de arbori, e ", 14, TXT), ("intrinsec bine calibrată", 14, INK, True), (".", 14, TXT)],
], size=14, gap=12)
figure(s, "calibration_logreg.png", 8.45, 2.0, 4.0, 4.2, caption="Diagrama de fiabilitate — regresie logistică")
footer(s, 12)
notes(s, "Metricile de până acum privesc ordinea piloților, nu calitatea probabilităților ca atare. "
         "Aici Random Forest domină din nou: ROC-AUC 0,94, cel mai mic log-loss și cea mai mică "
         "eroare de calibrare. În schimb, regresia logistică și arborele ordonează bine, dar sunt "
         "supraîncrezătoare. Le-am corectat cu scalare Platt și regresie izotonică pe validare; la "
         "logreg log-loss-ul scade de la 1,9 la 0,26. Important: calibrarea e monotonă, deci nu schimbă "
         "cine e prezis pe podium, doar face probabilitățile interpretabile.")

s = slide_new()
header(s, "Ce contează în predicție", "Importanța caracteristicilor", 13)
figure(s, "feature_importance_forest.png", LM, 2.0, 5.9, 4.4, caption="Random Forest — reducerea medie de impuritate")
points(s, 7.1, 2.1, 5.35, 4.3, [
    [("Poziția de start domină", 14.5, INK, True), (": ", 14.5, TXT), ("grid", 13, INK, True, False, MONO),
     (" și ", 14.5, TXT), ("quali_pos", 13, INK, True, False, MONO), (" pe primele locuri.", 14.5, TXT)],
    [("O caracteristică ", 14.5, TXT), ("nouă", 14.5, RED, True), (" urcă în top: ", 14.5, TXT),
     ("circuit_podium_rate_from_grid_pos", 11.5, INK, True, False, MONO), (".", 14.5, TXT)],
    [("Arborele unic concentrează ", 14.5, TXT), ("~61% într-o singură caracteristică", 14.5, INK, True), (" → varianță mare.", 14.5, TXT)],
    [("Pădurea distribuie importanța", 14.5, INK, True), (" pe multe caracteristici (cea mai mare < 10%) → robustețe.", 14.5, TXT)],
    [("Cele 15 caracteristici noi", 14.5, INK, True), (" ≈ 34% din importanța totală.", 14.5, TXT)],
], size=14.5, gap=12)
footer(s, 13)
notes(s, "Toate cele trei modele oferă interpretabilitate prin importanța caracteristicilor. Confirmă "
         "intuiția: poziția de start domină — grila și poziția în calificări sunt în vârf. Dar urcă și "
         "o caracteristică nouă, introdusă de mine: rata istorică de podium de pe poziția respectivă de "
         "start, care la regresia logistică e chiar prima. Un contrast revelator: arborele unic pune "
         "61% din importanță într-o singură caracteristică — de aici fragilitatea lui — în timp ce "
         "pădurea distribuie importanța, deci e robustă. Cele 15 caracteristici adăugate aduc aproximativ "
         "o treime din importanța totală.")

s = slide_new()
header(s, "Ablație & date in-sezon", "Două analize suplimentare", 14)
Text(tb(s, LM, 2.0, 5.6, 0.3)).add([("STUDIU DE ABLAȚIE", 12, MUTED, False, False, SEMI)], sa=0)
points(s, LM, 2.4, 5.6, 4.0, [
    "Eliminăm cea mai importantă caracteristică a fiecărui model și reantrenăm.",
    [("Rezultat contraintuitiv: performanța ", 13.5, TXT), ("nu scade", 13.5, INK, True), (" — ba chiar crește ușor pe unele metrici.", 13.5, TXT)],
    [("Cauza: ", 13.5, INK, True), ("redundanța", 13.5, TXT), (" — caracteristicile de start sunt corelate și se substituie reciproc.", 13.5, TXT)],
    [("Concluzie: ", 13.5, INK, True), ("importanța ≠ caracteristică de neînlocuit.", 13.5, TXT)],
], size=13.5, gap=11)
vline(s, 6.75, 2.05, 4.5)
Text(tb(s, 7.1, 2.0, 5.3, 0.3)).add([("EFECTUL DATELOR IN-SEZON (2025)", 12, MUTED, False, False, SEMI)], sa=0)
points(s, 7.1, 2.4, 5.3, 1.7, [
    "Antrenare incrementală cu cursele deja disputate din sezonul curent.",
    [("Efect real, dar ", 13.5, TXT), ("modest și neuniform", 13.5, INK, True), (".", 13.5, TXT)],
    [("Cel mai consecvent câștig: câștigător logreg ", 13.5, TXT), ("(Δ = +0,042)", 13.5, GREEN, True), (".", 13.5, TXT)],
], size=13.5, gap=8)
figure(s, "insezon_2025_exact_podium.png", 7.1, 4.25, 5.3, 2.1, frame=True)
footer(s, 14)
notes(s, "Două analize care întăresc rigoarea. Prima e un studiu de ablație: am eliminat din fiecare "
         "model cea mai importantă caracteristică și am reantrenat. Surprinzător, performanța nu scade, "
         "uneori crește puțin. Explicația e redundanța — caracteristicile legate de poziția de start "
         "sunt corelate, așa că, dacă scot una, celelalte îi preiau rolul. Concluzia metodologică: un "
         "scor mare de importanță nu înseamnă că acea caracteristică e de neînlocuit. A doua analiză e "
         "efectul datelor in-sezon: antrenarea incrementală aduce un câștig real, dar modest și neuniform.")

s = slide_new()
header(s, "Sinteză", "Concluzii", 15)
points(s, LM, 1.95, 7.5, 4.3, [
    [("Random Forest", 15, INK, True), (" — cea mai bună acuratețe (0,902), rată de podium (0,688) și podium exact (0,267).", 15, TXT)],
    [("Regresia logistică", 15, INK, True), (" — cea mai bună acuratețe a câștigătorului (0,541), de ~10× mai rapidă.", 15, TXT)],
    [("Arborele unic", 15, INK, True), (" — constant cel mai slab pe podium (varianță mare).", 15, TXT)],
    [("Sub dezechilibru, acuratețea e înșelătoare", 15, INK, True), (" → contează metricile la nivel de cursă.", 15, TXT)],
    [("Datele in-sezon", 15, INK, True), (" aduc un câștig real, dar modest și neuniform.", 15, TXT)],
], size=15, gap=12)
vline(s, 8.55, 2.0, 3.5)
Text(tb(s, 8.85, 1.95, 3.6, 0.3)).add([("CEL MAI BUN PER METRICĂ", 11.5, MUTED, False, False, SEMI)], sa=0)
byy = 2.42
for met, mod in [("Acuratețe", "Random Forest"), ("Rată podium", "Random Forest"),
                 ("Podium exact", "Random Forest"), ("Câștigător", "Regresie logistică")]:
    Text(tb(s, 8.85, byy, 3.6, 0.28)).add([(met, 12, MUTED)], sa=0)
    Text(tb(s, 8.85, byy + 0.24, 3.6, 0.35)).add([(mod, 14.5, INK, False, False, DISP)], sa=0)
    byy += 0.72
hairline(s, LM, 6.15, CW)
Text(tb(s, LM, 6.32, 11.5, 0.5)).add(
    [("✓  Toate cele 5 obiective propuse au fost îndeplinite integral.", 15.5, INK, False, False, DISP)], sa=0)
footer(s, 15)
notes(s, "În concluzie: Random Forest e modelul cu cea mai bună performanță medie, dominând trei din "
         "patru metrici, ceea ce confirmă avantajul ansamblurilor pe date tabelare. Regresia logistică, "
         "după optimizare, ia cea mai bună acuratețe a câștigătorului și e de zece ori mai rapidă. "
         "Arborele unic e cel mai slab pe podium, exact cum prezice teoria. Am arătat și de ce acuratețea "
         "e înșelătoare sub dezechilibru. Toate cele cinci obiective propuse au fost îndeplinite integral.")

s = slide_new()
header(s, "Direcții viitoare", "Dezvoltări ulterioare", 16)
points(s, LM, 2.0, 7.4, 4.2, [
    [("Algoritmi suplimentari", 15, INK, True), (": mașini cu vectori suport (SVM), gradient boosting (XGBoost).", 15, TXT)],
    [("Probabilități calibrate", 15, INK, True), (" integrate direct în agregarea clasamentului de campionat.", 15, TXT)],
    [("Caracteristici din telemetrie", 15, INK, True), (": ritm din antrenamente, modelare pe stint-uri a uzurii pneurilor.", 15, TXT)],
    [("Reformulare ca ", 15, TXT), ("learning-to-rank", 15, INK, True), (" pentru predicția poziției complete a fiecărui pilot.", 15, TXT)],
    [("Paralelizarea antrenării arborilor", 15, INK, True), (" — reducerea singurului dezavantaj al pădurii.", 15, TXT)],
], size=15, gap=13)
vline(s, 8.7, 2.1, 3.6)
Text(tb(s, 9.05, 3.15, 3.4, 1.0)).add([("Vă mulțumesc!", 32, INK, False, False, LIGHT)], sa=0)
Text(tb(s, 9.05, 4.05, 3.4, 0.6)).add([("Întrebări?", 19, RED, False, False, DISP)], sa=0)
footer(s, 16)
notes(s, "Pe scurt, soluția se poate extinde fără a-i schimba arhitectura: aș putea adăuga alți "
         "algoritmi, precum SVM sau gradient boosting, aș integra probabilitățile calibrate direct în "
         "clasament, aș folosi mai multe semnale din telemetrie și, pe termen lung, aș reformula "
         "problema ca learning-to-rank. Vă mulțumesc pentru atenție și aștept cu plăcere întrebările "
         "dumneavoastră.")

prs.save(OUT)
print("OK:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
